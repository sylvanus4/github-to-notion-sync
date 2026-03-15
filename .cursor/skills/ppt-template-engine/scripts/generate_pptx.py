#!/usr/bin/env python3
"""
Populate a PPTX template with content from a JSON spec.
Clones the template, then replaces placeholder markers ({{NAME}}) in each slide
with the corresponding values from the spec.

Usage:
  generate_pptx.py <template.pptx> <spec.json> <output.pptx>
"""

import json
import sys
import copy
from pptx import Presentation
from pptx.util import Pt


def _replace_text_in_shape(shape, replacements: dict):
    """Replace {{KEY}} markers in a shape's text frame."""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)
        for key, value in replacements.items():
            marker = f"{{{{{key}}}}}"
            if marker in full_text:
                if isinstance(value, list):
                    value = "\n".join(f"• {item}" for item in value)
                full_text = full_text.replace(marker, str(value))
        if paragraph.runs:
            paragraph.runs[0].text = full_text
            for run in paragraph.runs[1:]:
                run.text = ""


def generate(template_path: str, spec_path: str, output_path: str):
    with open(spec_path, "r") as f:
        spec = json.load(f)

    prs = Presentation(template_path)
    slides = list(prs.slides)

    # Build a layout-index to slide mapping from template
    layout_to_slide = {}
    with open(spec_path, "r") as f:
        spec_data = json.load(f)

    # Map template slides by index for duplication
    template_slide_count = len(slides)

    # Strategy: For each slide in spec, find the matching template slide
    # by layout_id, duplicate it, and fill placeholders.
    # We use a simple approach: map layout names to template slide indices
    # from the placeholder map.
    placeholder_map_path = spec_data.get("_placeholder_map_path")

    if placeholder_map_path:
        with open(placeholder_map_path, "r") as f:
            pmap = json.load(f)
        layout_index = {name: info["slide_index"] for name, info in pmap["layouts"].items()}
    else:
        # Fallback: use slide order as layout order
        layout_names = ["title-slide", "agenda", "two-column", "kpi-dashboard", "closing"]
        layout_index = {name: i for i, name in enumerate(layout_names) if i < template_slide_count}

    # For simplicity, we populate the existing template slides in order.
    # If the spec requests more slides of a type, we'd need XML-level duplication.
    # For v1, we support exactly the template slide count with placeholder replacement.

    spec_slides = spec_data.get("slides", [])

    # Replace placeholders in existing slides
    for i, slide_spec in enumerate(spec_slides):
        if i >= len(slides):
            break
        slide = slides[i]
        placeholders = slide_spec.get("placeholders", {})
        for shape in slide.shapes:
            _replace_text_in_shape(shape, placeholders)

    prs.save(output_path)
    print(json.dumps({
        "status": "success",
        "output": output_path,
        "slides_populated": min(len(spec_slides), len(slides)),
        "total_slides": len(slides),
    }, indent=2))


def main():
    if len(sys.argv) < 4:
        print("Usage: generate_pptx.py <template.pptx> <spec.json> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2], sys.argv[3])


if __name__ == "__main__":
    main()
