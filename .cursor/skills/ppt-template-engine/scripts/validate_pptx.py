#!/usr/bin/env python3
"""
Validate a populated PPTX against template rules.
Checks 8 compliance rules and outputs a JSON report.

Usage:
  validate_pptx.py <output.pptx> <placeholder-map.json>
"""

import json
import sys
import re
from pptx import Presentation


def validate(pptx_path: str, pmap_path: str) -> dict:
    with open(pmap_path, "r") as f:
        pmap = json.load(f)

    prs = Presentation(pptx_path)
    violations = []
    warnings = []
    stats = {"slides_checked": 0, "shapes_checked": 0}

    allowed_layouts = set(pmap["layouts"].keys())
    allowed_layout_indices = {info["slide_index"] for info in pmap["layouts"].values()}

    for i, slide in enumerate(prs.slides):
        stats["slides_checked"] += 1

        # Rule 1: Layout compliance - check slide index is within template range
        if i >= len(pmap["layouts"]):
            violations.append({
                "rule": "R1_LAYOUT_COMPLIANCE",
                "severity": "hard",
                "slide": i + 1,
                "message": f"Slide {i+1} exceeds template slide count ({len(pmap['layouts'])})",
            })
            continue

        # Rule 2: No rogue textboxes
        layout_name = list(pmap["layouts"].keys())[i] if i < len(pmap["layouts"]) else None
        if layout_name:
            allowed_names = set(pmap["layouts"][layout_name]["placeholders"].keys())
            allowed_names.update({"BG_RECT", "DIVIDER"})
            for shape in slide.shapes:
                stats["shapes_checked"] += 1
                if shape.name not in allowed_names and not shape.name.startswith("Slide Number"):
                    # Allow standard layout placeholders
                    if shape.name.startswith("Title") or shape.name.startswith("Content") or \
                       shape.name.startswith("Text") or shape.name.startswith("Date") or \
                       shape.name.startswith("Footer") or shape.name.startswith("Subtitle"):
                        continue
                    violations.append({
                        "rule": "R2_NO_ROGUE_SHAPES",
                        "severity": "hard",
                        "slide": i + 1,
                        "message": f"Unauthorized shape '{shape.name}' on slide {i+1}",
                    })

        # Rule 3: Placeholder fill check
        if layout_name:
            for ph_name, ph_info in pmap["layouts"][layout_name]["placeholders"].items():
                if not ph_info.get("required", False):
                    continue
                found = False
                for shape in slide.shapes:
                    if shape.name == ph_name and shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text and "{{" not in text:
                            found = True
                            break
                if not found:
                    marker_still_present = False
                    for shape in slide.shapes:
                        if shape.name == ph_name and shape.has_text_frame:
                            if "{{" in shape.text_frame.text:
                                marker_still_present = True
                    if marker_still_present:
                        violations.append({
                            "rule": "R3_PLACEHOLDER_UNFILLED",
                            "severity": "hard",
                            "slide": i + 1,
                            "message": f"Placeholder '{ph_name}' still contains marker template on slide {i+1}",
                        })

        # Rule 4: Theme/font/color preservation (check for unexpected font overrides)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and run.font.name not in (None, "Arial", "Pretendard", "Calibri"):
                            warnings.append({
                                "rule": "R4_THEME_PRESERVATION",
                                "severity": "soft",
                                "slide": i + 1,
                                "message": f"Non-standard font '{run.font.name}' in shape '{shape.name}'",
                            })

        # Rule 5: Title length check
        if layout_name:
            for ph_name, ph_info in pmap["layouts"][layout_name]["placeholders"].items():
                max_chars = ph_info.get("max_chars")
                if max_chars:
                    for shape in slide.shapes:
                        if shape.name == ph_name and shape.has_text_frame:
                            text = shape.text_frame.text
                            if len(text) > max_chars:
                                warnings.append({
                                    "rule": "R5_TEXT_LENGTH",
                                    "severity": "soft",
                                    "slide": i + 1,
                                    "message": f"'{ph_name}' text ({len(text)} chars) exceeds max ({max_chars})",
                                })

        # Rule 6: Bullet count check
        if layout_name:
            for ph_name, ph_info in pmap["layouts"][layout_name]["placeholders"].items():
                max_items = ph_info.get("max_items")
                if max_items and ph_info.get("type") == "bullet_list":
                    for shape in slide.shapes:
                        if shape.name == ph_name and shape.has_text_frame:
                            lines = [l for l in shape.text_frame.text.split("\n") if l.strip()]
                            if len(lines) > max_items:
                                warnings.append({
                                    "rule": "R6_BULLET_COUNT",
                                    "severity": "soft",
                                    "slide": i + 1,
                                    "message": f"'{ph_name}' has {len(lines)} items, max is {max_items}",
                                })

    # Rule 7: Check for remaining template markers anywhere
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                remaining = re.findall(r'\{\{\w+\}\}', text)
                if remaining:
                    violations.append({
                        "rule": "R7_RESIDUAL_MARKERS",
                        "severity": "hard",
                        "slide": i + 1,
                        "message": f"Unfilled markers {remaining} in shape '{shape.name}'",
                    })

    # Rule 8: Master integrity (basic check - slide count should match template)
    expected_slide_count = len(pmap["layouts"])
    if len(prs.slides) != expected_slide_count:
        warnings.append({
            "rule": "R8_MASTER_INTEGRITY",
            "severity": "soft",
            "message": f"Output has {len(prs.slides)} slides, template defines {expected_slide_count} layouts",
        })

    hard_fails = [v for v in violations if v["severity"] == "hard"]
    report = {
        "file": pptx_path,
        "passed": len(hard_fails) == 0,
        "hard_violations": len(hard_fails),
        "soft_warnings": len(warnings),
        "violations": violations,
        "warnings": warnings,
        "stats": stats,
    }
    return report


def main():
    if len(sys.argv) < 3:
        print("Usage: validate_pptx.py <output.pptx> <placeholder-map.json>", file=sys.stderr)
        sys.exit(1)

    report = validate(sys.argv[1], sys.argv[2])
    print(json.dumps(report, indent=2, ensure_ascii=False))
    sys.exit(0 if report["passed"] else 1)


if __name__ == "__main__":
    main()
