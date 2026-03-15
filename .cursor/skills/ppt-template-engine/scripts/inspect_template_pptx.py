#!/usr/bin/env python3
"""
Inspect a PPTX template: extract slide layouts, shape names, placeholder
markers ({{...}}), and theme colors. Outputs JSON to stdout.
"""

import json
import sys
from pptx import Presentation
from pptx.util import Emu


def inspect(pptx_path: str) -> dict:
    prs = Presentation(pptx_path)
    result = {
        "file": pptx_path,
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "slide_count": len(prs.slides),
        "slide_layouts": [],
        "slides": [],
    }

    for layout in prs.slide_layouts:
        layout_info = {
            "name": layout.name,
            "placeholders": [
                {"idx": ph.placeholder_format.idx, "name": ph.name, "type": str(ph.placeholder_format.type)}
                for ph in layout.placeholders
            ],
        }
        result["slide_layouts"].append(layout_info)

    for i, slide in enumerate(prs.slides):
        shapes_info = []
        for shape in slide.shapes:
            info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": round(shape.left / 914400, 2) if shape.left else 0,
                "top": round(shape.top / 914400, 2) if shape.top else 0,
                "width": round(shape.width / 914400, 2) if shape.width else 0,
                "height": round(shape.height / 914400, 2) if shape.height else 0,
            }
            if shape.has_text_frame:
                info["text"] = shape.text_frame.text[:200]
                info["is_placeholder_marker"] = "{{" in shape.text_frame.text
            shapes_info.append(info)

        slide_info = {
            "index": i + 1,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "unknown",
            "shapes": shapes_info,
            "placeholder_names": [
                s["name"] for s in shapes_info
                if s["name"] not in ("BG_RECT", "DIVIDER") and s["name"].isupper()
            ],
        }
        result["slides"].append(slide_info)

    return result


def main():
    if len(sys.argv) < 2:
        print("Usage: inspect_template_pptx.py <file.pptx>", file=sys.stderr)
        sys.exit(1)

    data = inspect(sys.argv[1])
    print(json.dumps(data, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
