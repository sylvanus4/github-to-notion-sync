#!/usr/bin/env python3
"""Create ThakiCloud proposal PPTX template with named placeholders."""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

THAKI_BLUE = RGBColor(0x1A, 0x56, 0xDB)
THAKI_DARK = RGBColor(0x1E, 0x29, 0x3B)
THAKI_LIGHT = RGBColor(0xF0, 0xF4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=THAKI_DARK, alignment=PP_ALIGN.LEFT, name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if name:
        txBox.name = name
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_bg_rect(slide, prs, color=THAKI_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.name = "BG_RECT"


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg_rect(slide, prs, THAKI_BLUE)
    _add_textbox(slide, 1.0, 2.0, 8.0, 1.0, "{{TITLE}}", 36, True, WHITE, PP_ALIGN.CENTER, "TITLE")
    _add_textbox(slide, 1.0, 3.2, 8.0, 0.6, "{{SUBTITLE}}", 20, False, THAKI_LIGHT, PP_ALIGN.CENTER, "SUBTITLE")
    _add_textbox(slide, 1.0, 4.5, 8.0, 0.5, "{{AUTHOR}}", 16, False, THAKI_LIGHT, PP_ALIGN.CENTER, "AUTHOR")


def create_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.3, 9.0, 0.7, "{{SECTION_TITLE}}", 28, True, THAKI_BLUE, PP_ALIGN.LEFT, "SECTION_TITLE")
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9.0), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = THAKI_BLUE
    shape.line.fill.background()
    shape.name = "DIVIDER"
    _add_textbox(slide, 0.8, 1.5, 8.5, 4.5, "{{AGENDA_ITEMS}}", 18, False, THAKI_DARK, PP_ALIGN.LEFT, "AGENDA_ITEMS")


def create_two_column_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.3, 4.0, 0.6, "{{LEFT_TITLE}}", 22, True, THAKI_BLUE, PP_ALIGN.LEFT, "LEFT_TITLE")
    _add_textbox(slide, 0.5, 1.1, 4.2, 4.5, "{{LEFT_BODY}}", 14, False, THAKI_DARK, PP_ALIGN.LEFT, "LEFT_BODY")
    _add_textbox(slide, 5.3, 0.3, 4.0, 0.6, "{{RIGHT_TITLE}}", 22, True, THAKI_BLUE, PP_ALIGN.LEFT, "RIGHT_TITLE")
    _add_textbox(slide, 5.3, 1.1, 4.2, 4.5, "{{RIGHT_BODY}}", 14, False, THAKI_DARK, PP_ALIGN.LEFT, "RIGHT_BODY")


def create_kpi_dashboard_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.3, 9.0, 0.7, "{{KPI_TITLE}}", 28, True, THAKI_BLUE, PP_ALIGN.LEFT, "KPI_TITLE")
    _add_textbox(slide, 0.5, 1.3, 9.0, 4.5, "{{KPI_TABLE}}", 14, False, THAKI_DARK, PP_ALIGN.LEFT, "KPI_TABLE")


def create_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, prs, THAKI_BLUE)
    _add_textbox(slide, 1.0, 2.5, 8.0, 1.0, "{{CLOSING_MESSAGE}}", 32, True, WHITE, PP_ALIGN.CENTER, "CLOSING_MESSAGE")
    _add_textbox(slide, 1.0, 4.0, 8.0, 0.5, "{{CONTACT}}", 16, False, THAKI_LIGHT, PP_ALIGN.CENTER, "CONTACT")


def main():
    output_path = sys.argv[1] if len(sys.argv) > 1 else "thaki-proposal-v1.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_agenda_slide(prs)
    create_two_column_slide(prs)
    create_kpi_dashboard_slide(prs)
    create_closing_slide(prs)

    prs.save(output_path)
    print(f"Template created: {output_path}")
    print(f"Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        placeholders = [s.name for s in slide.shapes if s.name.startswith("{{") or s.name.isupper()]
        print(f"  Slide {i+1}: {placeholders}")


if __name__ == "__main__":
    main()
